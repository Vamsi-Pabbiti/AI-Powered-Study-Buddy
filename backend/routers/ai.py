from fastapi import APIRouter, Depends, HTTPException, UploadFile, File, Form
from models.schemas import ExplainRequest, SummarizeRequest, QuizRequest, FlashcardRequest
from utils.auth_utils import get_current_user
from utils.database import get_db
from datetime import datetime
from pathlib import Path
from dotenv import load_dotenv
from pptx import Presentation
from groq import Groq
import os, json
import fitz  # pymupdf
from docx import Document
import io

load_dotenv(Path(__file__).resolve().parents[1] / ".env")
router = APIRouter()
client = Groq(api_key=os.getenv("GROQ_API_KEY"))

def ask_ai(prompt: str) -> str:
    try:
        response = client.chat.completions.create(
            model="llama-3.3-70b-versatile",
            messages=[{"role": "user", "content": prompt}],
            max_tokens=1500
        )
        return response.choices[0].message.content
    except Exception as exc:
        raise HTTPException(status_code=502, detail=f"Groq API failed: {exc}")

def parse_json_array(text: str):
    cleaned = text.strip().replace("```json", "").replace("```", "").strip()
    try:
        parsed = json.loads(cleaned)
    except json.JSONDecodeError:
        start = cleaned.find("[")
        end = cleaned.rfind("]")
        if start == -1 or end == -1 or end <= start:
            raise ValueError("AI response was not valid JSON")
        parsed = json.loads(cleaned[start:end + 1])
    if not isinstance(parsed, list):
        raise ValueError("AI response was not a JSON array")
    return parsed

async def save_history(db, user_id, feature, input_text, output_text):
    await db.history.insert_one({
        "user_id": user_id, "feature": feature,
        "input": input_text, "output": output_text,
        "created_at": datetime.utcnow()
    })

@router.post("/explain")
async def explain_topic(req: ExplainRequest, user_id: str = Depends(get_current_user)):
    db = get_db()
    level_map = {
        "simple": "Explain like I'm 10 years old using simple words and fun analogies.",
        "medium": "Explain clearly for a high school student with examples.",
        "advanced": "Explain in depth for a college student, including technical details."
    }
    prompt = f"{level_map[req.level]}\n\nTopic: {req.topic}\n\nGive a well-structured explanation with key points."
    result = ask_ai(prompt)
    await save_history(db, user_id, "explain", req.topic, result)
    return {"result": result}

def extract_text_from_file(file_bytes: bytes, filename: str) -> str:
    ext = filename.lower().split(".")[-1]
    if ext == "pdf":
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        return "\n".join(page.get_text() for page in doc)
    elif ext == "docx":
        doc = Document(io.BytesIO(file_bytes))
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    elif ext == "txt":
        return file_bytes.decode("utf-8", errors="ignore")
    elif ext == "pptx":
        prs = Presentation(io.BytesIO(file_bytes))
        slides_text = []
        for i, slide in enumerate(prs.slides, start=1):
            slide_lines = [f"--- Slide {i} ---"]
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    slide_lines.append(shape.text_frame.text.strip())
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = " | ".join(cell.text.strip() for cell in row.cells)
                        if row_text.strip(" |"):
                            slide_lines.append(row_text)
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
                slide_lines.append(f"Notes: {slide.notes_slide.notes_text_frame.text.strip()}")
            if len(slide_lines) > 1:
                slides_text.append("\n".join(slide_lines))
        return "\n\n".join(slides_text)
    else:
        raise HTTPException(status_code=400, detail="Unsupported file type. Use PDF, DOCX, PPTX, or TXT.")

@router.post("/summarize")
async def summarize_notes(
    style: str = Form("bullets"),
    notes: str = Form(""),
    file: UploadFile = File(None),
    user_id: str = Depends(get_current_user)
):
    db = get_db()

    # Get text from file or plain text input
    if file and file.filename:
        file_bytes = await file.read()
        text = extract_text_from_file(file_bytes, file.filename)
        source_label = file.filename
    elif notes.strip():
        text = notes
        source_label = notes[:100] + "..."
    else:
        raise HTTPException(status_code=400, detail="Provide either notes text or upload a file.")

    if len(text.strip()) < 20:
        raise HTTPException(status_code=400, detail="Extracted text is too short to summarize.")

    style_map = {
        "bullets": "Summarize as clear bullet points.",
        "paragraph": "Summarize as a concise paragraph.",
        "key_points": "Extract the top 5-7 key points only."
    }
    prompt = f"{style_map.get(style, style_map['bullets'])}\n\nNotes:\n{text[:6000]}"
    result = ask_ai(prompt)
    await save_history(db, user_id, "summarize", source_label, result)
    return {"result": result}

@router.post("/quiz")
async def generate_quiz(req: QuizRequest, user_id: str = Depends(get_current_user)):
    db = get_db()
    prompt = f"""Generate {req.num_questions} multiple choice questions on "{req.topic}" at {req.difficulty} difficulty.

Return ONLY a JSON array in this exact format (no markdown, no extra text):
[
  {{
    "question": "Question text here?",
    "options": ["A) ...", "B) ...", "C) ...", "D) ..."],
    "answer": "A) ...",
    "explanation": "Brief explanation why this is correct."
  }}
]"""
    text = ask_ai(prompt)
    try:
        questions = parse_json_array(text)
    except (json.JSONDecodeError, ValueError) as exc:
        raise HTTPException(status_code=502, detail=f"Quiz JSON parse failed: {exc}")
    await save_history(db, user_id, "quiz", req.topic, f"{req.num_questions} questions generated")
    return {"questions": questions}

@router.post("/flashcards")
async def generate_flashcards(req: FlashcardRequest, user_id: str = Depends(get_current_user)):
    db = get_db()
    prompt = f"""Create {req.num_cards} flashcards for studying "{req.topic}".

Return ONLY a JSON array in this exact format (no markdown, no extra text):
[
  {{
    "front": "Question or term",
    "back": "Answer or definition"
  }}
]"""
    text = ask_ai(prompt)
    try:
        cards = parse_json_array(text)
    except (json.JSONDecodeError, ValueError) as exc:
        raise HTTPException(status_code=502, detail=f"Flashcard JSON parse failed: {exc}")
    await save_history(db, user_id, "flashcard", req.topic, f"{req.num_cards} flashcards generated")
    return {"flashcards": cards}